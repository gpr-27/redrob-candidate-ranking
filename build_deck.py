#!/usr/bin/env python3
"""Fill the Redrob idea-submission PPTX template with our approach content."""

import copy
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

DARK = RGBColor(0x20, 0x27, 0x29)
FONT = "Manrope SemiBold"

SRC = "Idea Submission Template _ Redrob.pptx"
OUT = "Redrob_Idea_Submission_Praneeth.pptx"

REPO = "https://github.com/gpr-27/redrob-candidate-ranking"


def fill_body(shape, lines, size=11):
    """Replace a text frame's content with the given (text, bold, indent) lines."""
    tf = shape.text_frame
    tf.word_wrap = True
    # keep first paragraph, clear the rest
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, (text, bold, indent) in enumerate(lines):
        p = first if i == 0 else tf.add_paragraph()
        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = DARK


prs = Presentation(SRC)
slides = list(prs.slides)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
title_map = {
    "Team Name :": "Team Name :  Team Praneeth",
    "Team Leader Name :": "Team Leader Name :  Praneeth",
    "Problem Statement :": "Problem Statement :  Data & AI Challenge — Intelligent Candidate Discovery & Ranking",
}
for shape in slides[0].shapes:
    if shape.has_text_frame:
        cur = shape.text_frame.paragraphs[0].text.strip()
        if cur in title_map:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = title_map[cur]
                    break
                break

# ── Slide 2: Solution Overview ──────────────────────────────────────────────
body2 = [
    ("A deterministic, rule-based ranking engine that scores all 100,000 candidates against the Senior AI/ML Engineer JD and outputs the top-100 in ~10 seconds on CPU — pure Python standard library, no network, no GPU, no LLM calls.", False, 0),
    ("", False, 0),
    ("What differentiates it from traditional candidate matching:", True, 0),
    ("Encodes the JD's actual hiring rubric as 7 weighted scoring dimensions instead of embedding-cosine keyword similarity.", False, 1),
    ("Values skill QUALITY over quantity: proficiency x months-of-use, so an 'expert with 0 months' scores zero.", False, 1),
    ("Weighs hireability, not just fit: platform recency, recruiter response rate, and notice period modify the score.", False, 1),
    ("Built-in impossible-profile (honeypot) detection zeroes out fabricated candidates.", False, 1),
    ("Fully deterministic and reproducible: same input, byte-identical output, every run.", False, 1),
]
fill_body(slides[1].shapes[2], body2)

# ── Slide 3: JD Understanding ───────────────────────────────────────────────
body3 = [
    ("Key requirements extracted from the JD:", True, 0),
    ("5-9 years experience (ideal 6-8), with 4-5 in applied ML at product companies; production embeddings/retrieval + vector DB experience; strong Python; ranking evaluation (NDCG, MRR); India, Pune/Noida preferred; sub-30-day notice; active on the Redrob platform.", False, 1),
    ("Explicit anti-signals: consulting-only careers, CV/speech-only backgrounds, title-chasers, framework enthusiasts.", False, 1),
    ("", False, 0),
    ("Most important candidate signals (in scoring order):", True, 0),
    ("Current + historical job titles (25%) -> a Marketing Manager with AI keywords is not a fit.", False, 1),
    ("Career history text evidence (25%) -> did they actually build ranking / search / recommendation systems?", False, 1),
    ("Skills with proficiency x duration (20%), behavioral engagement (15%), experience band (8%), location (4%), education (3%).", False, 1),
    ("", False, 0),
    ("Beyond keyword matching: title tiering, career-description evidence mining, duration-weighted skill trust, and cross-field consistency checks catch what keyword counts cannot.", False, 0),
]
fill_body(slides[2].shapes[2], body3, size=10)

# ── Slide 4: Ranking Methodology ────────────────────────────────────────────
body4 = [
    ("Score = (0.25 Title + 0.25 Career + 0.20 Skills + 0.15 Behavioral + 0.08 Experience + 0.04 Location + 0.03 Education) x HoneypotMultiplier", True, 0),
    ("", False, 0),
    ("Title: 5-tier classification (Staff ML / Search / RecSys engineer = 1.0 ... non-technical = 0.0), blending current (70%) and best historical (30%) title.", False, 1),
    ("Career: product-vs-consulting company weighting (all-consulting = 0.25x penalty), industry relevance, JD-keyword evidence in role descriptions, tenure health (job-hopper penalty).", False, 1),
    ("Skills: must-have JD skills weighted by proficiency x duration_months; wrong-domain (CV/speech) and keyword-stuffing penalties; Redrob assessment scores.", False, 1),
    ("Behavioral: last-active recency, recruiter response rate & speed, open-to-work + notice period, market demand, interview/offer reliability.", False, 1),
    ("", False, 0),
    ("All 100K candidates are streamed and scored; sorted by (-score, candidate_id) so equal scores tie-break deterministically; top-100 selected with per-candidate reasoning.", False, 0),
    ("No black-box model: every heuristic is a JD requirement made executable - fully explainable and defensible.", False, 0),
]
fill_body(slides[3].shapes[2], body4, size=10)

# ── Slide 5: Explainability & Data Validation ───────────────────────────────
body5 = [
    ("Explainability:", True, 0),
    ("Every top-100 row carries reasoning assembled ONLY from facts present in that candidate's profile: title, company, years, actual matching skills, education tier, behavioral signals - plus honest concerns (long notice, low response rate, non-India location).", False, 1),
    ("Hallucination is prevented by construction: the generator can only reference profile fields, never generate free text. Verified programmatically - 100/100 reasoning strings match the underlying profiles exactly, all unique.", False, 1),
    ("", False, 0),
    ("Data validation / suspicious-profile handling (honeypot detector):", True, 0),
    ("Employment starting before the company was founded (checked against real founding years: Krutrim 2023, Sarvam AI 2023, Glance 2019, ...).", False, 1),
    ("Stated years_of_experience exceeding the entire career timeline; summary text contradicting the experience field.", False, 1),
    ("'Expert' skills with 0 months of use; mass all-expert skill lists; non-technical titles with deep ML skill stacks.", False, 1),
    ("Suspicion maps to a 1.0 -> 0.0 multiplier. Final top-100 audit: 0 impossible profiles (disqualification threshold is 10).", False, 1),
]
fill_body(slides[4].shapes[2], body5, size=10)

# ── Slide 6: End-to-End Workflow ────────────────────────────────────────────
body6 = [
    ("1.  Stream candidates.jsonl line-by-line (constant memory, no full-file load).", False, 0),
    ("2.  For each candidate, compute the 7 component scores (title, career, skills, behavioral, experience, location, education).", False, 0),
    ("3.  Combine via the weighted sum; apply the honeypot multiplier (1.0 clean -> 0.0 impossible).", False, 0),
    ("4.  Global sort by (-score, candidate_id ascending) - deterministic tie-breaking per the spec.", False, 0),
    ("5.  Select top-100; round scores to 4 decimals; re-sort on rounded scores to keep the validator's tie-break rule intact.", False, 0),
    ("6.  Generate fact-grounded reasoning for each of the 100 candidates.", False, 0),
    ("7.  Write submission.csv (validated by the official validate_submission.py).", False, 0),
    ("", False, 0),
    ("Single command:  python rank.py --candidates ./candidates.jsonl --out ./submission.csv", True, 0),
    ("End-to-end: ~10 seconds for 100,000 candidates on a CPU-only machine.", False, 0),
]
fill_body(slides[5].shapes[2], body6)

# ── Slide 7: System Architecture (no body box in template — add one) ────────
slide7 = slides[6]
box = slide7.shapes.add_textbox(Inches(0.41), Inches(1.40), Inches(9.32), Inches(3.5))
body7 = [
    ("candidates.jsonl (100K)", True, 0),
    ("        |  stream, one candidate at a time", False, 0),
    ("        v", False, 0),
    ("rank.py orchestrator", True, 0),
    ("        |-- scoring/title_scorer.py         (25%)  5-tier title classification", False, 0),
    ("        |-- scoring/career_scorer.py        (25%)  companies, industries, JD-evidence, tenure", False, 0),
    ("        |-- scoring/skills_scorer.py        (20%)  proficiency x duration, stuffing penalty", False, 0),
    ("        |-- scoring/behavioral_scorer.py    (15%)  recency, response rate, notice period", False, 0),
    ("        |-- scoring/experience_scorer.py    ( 8%)  5-9 yr bell curve", False, 0),
    ("        |-- scoring/location_scorer.py      ( 4%)  India / Pune / Noida / relocation", False, 0),
    ("        |-- scoring/education_scorer.py     ( 3%)  field, tier, degree", False, 0),
    ("        |-- scoring/honeypot_detector.py    (x multiplier)  impossible-profile filter", False, 0),
    ("        |-- scoring/reasoning_generator.py  fact-grounded explanations", False, 0),
    ("        v", False, 0),
    ("sort (-score, candidate_id)  ->  top-100  ->  submission.csv", True, 0),
]
fill_body(box, body7, size=9)
for p in box.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = "Courier New"

# ── Slide 8: Results & Performance ──────────────────────────────────────────
body8 = [
    ("Ranking quality (final top-100):", True, 0),
    ("100/100 hold ML/AI/search/recommendation titles; 94/100 based in India; 86/100 inside the 5-9 year band.", False, 1),
    ("0 honeypots / impossible profiles (independently re-audited after ranking; disqualification threshold is 10).", False, 1),
    ("0 consulting-only careers; 0 keyword stuffers (lowest avg expert-skill duration in top-100: 55 months).", False, 1),
    ("The only two candidates in the pool matching ALL SEVEN of the JD's ideal-profile criteria rank #3 and #8.", False, 1),
    ("100/100 unique, fact-verified reasoning strings.", False, 1),
    ("", False, 0),
    ("Compute constraints (limit vs measured):", True, 0),
    ("Runtime: 5 min allowed -> ~10 s actual.   Memory: 16 GB allowed -> ~1.8 GB peak.", False, 1),
    ("CPU-only, zero network calls, zero GPU - verified in a fresh clean-room clone that reproduced the CSV byte-for-byte.", False, 1),
    ("Passes the official validate_submission.py with no issues.", False, 1),
]
fill_body(slides[7].shapes[2], body8, size=10)

# ── Slide 9: Technologies Used ──────────────────────────────────────────────
body9 = [
    ("Python 3.11+ standard library only - json, csv, datetime, re, argparse, pathlib.", True, 0),
    ("", False, 0),
    ("Why this stack:", True, 0),
    ("Determinism & reproducibility: no model weights, no random seeds, no API drift. Stage-3 reproduction cannot fail.", False, 1),
    ("Speed: streaming + pure-Python heuristics rank 100K candidates in ~10 s - 30x under the 5-minute budget.", False, 1),
    ("Zero dependencies: requirements.txt is intentionally empty; nothing to install, nothing to break in the sandbox.", False, 1),
    ("Explainability: every scoring rule maps 1:1 to a sentence in the JD - defensible line-by-line at interview.", False, 1),
    ("", False, 0),
    ("Deliberately avoided: per-candidate LLM calls (cannot fit the compute budget and hallucinate), embedding similarity as the primary signal (falls for keyword stuffers - the exact trap the dataset sets).", False, 0),
    ("AI tools declared: Claude (architecture discussion, scoring design, implementation & review). No candidate data was sent to any LLM.", False, 0),
]
fill_body(slides[8].shapes[2], body9, size=10)

# ── Slide 10: Submission Assets ─────────────────────────────────────────────
body10 = [
    ("GitHub repository (public):", True, 0),
    (REPO, False, 1),
    ("", False, 0),
    ("Live Sandbox Demo (verified):", True, 0),
    ("https://huggingface.co/spaces/praneethg27/redrob-ranker", False, 1),
    ("", False, 0),
    ("Reproduce command (from repo README):", True, 0),
    ("python rank.py --candidates ./candidates.jsonl --out ./submission.csv", False, 1),
    ("", False, 0),
    ("Repo contents: full source (rank.py + 9 scoring modules), README with methodology and setup, requirements.txt, submission_metadata.yaml, submission.csv, official validator.", False, 0),
    ("", False, 0),
    ("Ranked output: submission.csv / submission.xlsx (top-100 with rank, score, and per-candidate reasoning).", False, 0),
    ("Sandbox accepts custom uploads and produces valid submission.csv.", False, 0),
]
fill_body(slides[9].shapes[2], body10)

prs.save(OUT)
print(f"Saved {OUT}")
